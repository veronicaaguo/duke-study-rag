"""
src/ingestion/loader.py

Loads course documents from disk into a unified format.
Supports: PDF, PPTX, DOCX, TXT, MD

When use_vision=True (default), PDF and PPTX files are processed with
GPT-4o vision for image-heavy pages — see vision_loader.py.
Text-only files (DOCX, TXT, MD) are never sent to the vision API.
"""

from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Optional
from docx import Document
from loguru import logger


@dataclass
class RawDocument:
    """A single loaded document before chunking."""
    content: str
    source: str          # original file path
    course: str          # e.g. "CS372"
    doc_type: str        # "pdf", "pdf_vision", "pptx", "pptx_vision", "docx", "txt"
    page_count: int = 0
    metadata: dict = field(default_factory=dict)


def load_pdf(path: Path, course: str) -> RawDocument:
    """Text-only PDF loader (no vision). Used as fallback or when vision disabled."""
    import fitz
    doc = fitz.open(str(path))
    pages = []
    for page_num, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(f"[Page {page_num + 1}]\n{text}")
    content = "\n\n".join(pages)
    return RawDocument(
        content=content,
        source=str(path),
        course=course,
        doc_type="pdf",
        page_count=len(doc),
        metadata={"filename": path.name}
    )


def load_pptx(path: Path, course: str) -> RawDocument:
    """Text-only PPTX loader. Used as fallback or when vision disabled."""
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
    content = "\n\n".join(slides)
    return RawDocument(
        content=content,
        source=str(path),
        course=course,
        doc_type="pptx",
        page_count=len(prs.slides),
        metadata={"filename": path.name}
    )


def load_docx(path: Path, course: str) -> RawDocument:
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    content = "\n\n".join(paragraphs)
    return RawDocument(
        content=content,
        source=str(path),
        course=course,
        doc_type="docx",
        page_count=0,
        metadata={"filename": path.name}
    )


def load_txt(path: Path, course: str) -> RawDocument:
    content = path.read_text(encoding="utf-8", errors="ignore")
    return RawDocument(
        content=content,
        source=str(path),
        course=course,
        doc_type="txt",
        page_count=0,
        metadata={"filename": path.name}
    )


# Text-only loaders (always available)
TEXT_LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".ppt": load_pptx,
    ".docx": load_docx,
    ".doc": load_docx,
    ".txt": load_txt,
    ".md": load_txt,
}


def _load_single(path: Path, course: str, use_vision: bool) -> Optional[RawDocument]:
    """
    Load one file. Routes PDF/PPTX through GPT-4o vision loader when enabled.
    Returns None if the file yields no content.
    """
    suffix = path.suffix.lower()

    if use_vision and suffix in (".pdf",):
        from src.ingestion.vision_loader import load_pdf_vision
        return load_pdf_vision(path, course, use_vision=True)

    if use_vision and suffix in (".pptx", ".ppt"):
        from src.ingestion.vision_loader import load_pptx_vision
        return load_pptx_vision(path, course, use_vision=True)

    if suffix in TEXT_LOADERS:
        return TEXT_LOADERS[suffix](path, course)

    return None


def load_directory(
    input_dir: Path,
    course: str,
    use_vision: bool = True,
) -> List[RawDocument]:
    """
    Load all supported documents from a directory (recursive).

    Args:
        input_dir:   Directory containing course files.
        course:      Course name tag, e.g. "CS372".
        use_vision:  If True, send image-heavy PDF/PPTX pages to GPT-4o.
                     Set False to skip vision API calls (faster, cheaper, text-only).
    """
    docs = []
    supported = set(TEXT_LOADERS.keys())

    for path in sorted(input_dir.rglob("*")):
        if path.suffix.lower() not in supported:
            continue
        try:
            doc = _load_single(path, course, use_vision=use_vision)
            if doc and doc.content.strip():
                docs.append(doc)
                vision_tag = " [+vision]" if "vision" in doc.doc_type else ""
                logger.info(
                    f"Loaded {path.name}{vision_tag} "
                    f"({doc.page_count} pages, {len(doc.content)} chars, "
                    f"vision_calls={doc.metadata.get('vision_calls', 0)})"
                )
            else:
                logger.warning(f"Skipped {path.name} — no extractable content")
        except Exception as e:
            logger.error(f"Failed to load {path.name}: {e}")

    logger.info(f"Loaded {len(docs)} documents from {input_dir} (vision={'on' if use_vision else 'off'})")
    return docs
